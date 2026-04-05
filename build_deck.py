#!/usr/bin/env python3
"""
BET Deck Builder — PPTX generation engine.

Reads a JSON outline and produces a CD-compliant .pptx using the
BET Ergebnis-Master template.

Usage:
    python build_deck.py --template path/to/Master.pptx --outline outline.json --output deck.pptx

Outline JSON format:
{
    "title": "Presentation Title",
    "subtitle": "Optional subtitle",
    "footer": "Project Name | Month Year",
    "slides": [
        {
            "layout": 2,
            "title": "Single-line action title",
            "body": ["Bullet 1", "Bullet 2", "→ Fazit line"]
        },
        {
            "layout": 6,
            "title": "Four equal items",
            "kacheln": [
                {"title": "Item 1", "body": "Description"},
                {"title": "Item 2", "body": "Description"},
                {"title": "Item 3", "body": "Description"},
                {"title": "Item 4", "body": "Description"}
            ]
        }
    ]
}
"""

import argparse
import json
import copy
import datetime
import sys

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from lxml import etree

NS_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'
NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
nsmap = {'p': NS_P, 'a': NS_A}

BLACK = RGBColor(0x00, 0x00, 0x00)

# Kachel index mapping for Layout 6
KACHEL_MAP = [
    (24, 22),   # top-left:     title_idx, body_idx
    (28, 27),   # top-right
    (26, 25),   # bottom-left
    (30, 29),   # bottom-right
]


def _remove_bullet_and_indent(paragraph):
    pPr = paragraph._p.get_or_add_pPr()
    etree.SubElement(pPr, f'{{{NS_A}}}buNone')
    pPr.set('marL', '0')
    pPr.set('indent', '0')


def set_text_inherit(placeholder, text, bold=False):
    tf = placeholder.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.color.rgb = BLACK
    if bold:
        run.font.bold = True


def add_bullets_inherit(placeholder, items):
    tf = placeholder.text_frame
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        is_fazit = item.startswith("→ ")
        if is_fazit:
            _remove_bullet_and_indent(p)
            p.space_before = Pt(6)
            run = p.add_run()
            run.text = item
            run.font.bold = True
            run.font.color.rgb = BLACK
        else:
            run = p.add_run()
            run.text = item
            run.font.color.rgb = BLACK
        p.space_after = Pt(2)
        p.level = 0


def set_kachel(slide, title_idx, body_idx, title_text, body_text):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == title_idx:
            set_text_inherit(ph, title_text, bold=True)
        elif ph.placeholder_format.idx == body_idx:
            tf = ph.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            _remove_bullet_and_indent(p)
            run = p.add_run()
            run.text = body_text
            run.font.color.rgb = BLACK


def clone_footer_and_slidenum(slide, layout, slide_number, footer_text):
    spTree = slide.element.find(f'{{{NS_P}}}cSld/{{{NS_P}}}spTree')
    for sp in layout.element.findall('.//p:sp', nsmap):
        ph_elem = sp.find('.//p:nvSpPr/p:nvPr/p:ph', nsmap)
        if ph_elem is not None:
            idx = ph_elem.get('idx', '0')
            if idx in ['10', '11']:
                new_sp = copy.deepcopy(sp)
                max_id = max(
                    (int(el.get('id', '0'))
                     for el in spTree.findall('.//{%s}cNvPr' % NS_A)),
                    default=100
                )
                cNvPr = new_sp.find('.//{%s}cNvPr' % NS_A)
                if cNvPr is not None:
                    cNvPr.set('id', str(max_id + 1))
                txBody = new_sp.find('.//{%s}txBody' % NS_A)
                if txBody is not None:
                    for p in txBody.findall(f'{{{NS_A}}}p'):
                        txBody.remove(p)
                    p_elem = etree.SubElement(txBody, f'{{{NS_A}}}p')
                    r_elem = etree.SubElement(p_elem, f'{{{NS_A}}}r')
                    rPr = etree.SubElement(r_elem, f'{{{NS_A}}}rPr')
                    rPr.set('lang', 'de-DE')
                    rPr.set('dirty', '0')
                    t_elem = etree.SubElement(r_elem, f'{{{NS_A}}}t')
                    t_elem.text = footer_text if idx == '10' else str(slide_number)
                spTree.append(new_sp)


def remove_existing_slides(prs):
    for _ in range(len(prs.slides)):
        sldId = prs.slides._sldIdLst[0]
        rId = sldId.get(
            '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
        )
        if rId is None:
            rId = sldId.get('r:id')
        if rId:
            try:
                prs.part.drop_rel(rId)
            except KeyError:
                pass
        prs.slides._sldIdLst.remove(sldId)


def build_title_slide(prs, layouts, outline):
    slide = prs.slides.add_slide(layouts[0])
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            set_text_inherit(ph, outline['title'], bold=True)
        elif ph.placeholder_format.idx == 15:
            subtitle = outline.get('subtitle', '')
            if subtitle:
                set_text_inherit(ph, subtitle)
        elif ph.placeholder_format.idx == 13:
            set_text_inherit(ph, datetime.date.today().strftime("%d.%m.%Y"))
    return slide


def build_content_slide(prs, layouts, slide_def, slide_num, footer_text):
    layout_idx = slide_def['layout']
    layout = layouts[layout_idx]
    slide = prs.slides.add_slide(layout)

    # Set title
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            set_text_inherit(ph, slide_def['title'], bold=True)

    # Handle body text (layouts 2, 3, 4, 5)
    if 'body' in slide_def:
        body_idx = 15  # default
        if layout_idx == 7:
            body_idx = 21
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == body_idx:
                add_bullets_inherit(ph, slide_def['body'])

    # Handle kacheln (layout 6)
    if 'kacheln' in slide_def:
        kacheln = slide_def['kacheln']
        for k, (t_idx, b_idx) in zip(kacheln, KACHEL_MAP):
            set_kachel(slide, t_idx, b_idx, k['title'], k['body'])

    # Clone footer and slide number
    clone_footer_and_slidenum(slide, layout, slide_num, footer_text)

    return slide


def build(template_path, outline, output_path):
    prs = Presentation(template_path)
    remove_existing_slides(prs)
    layouts = list(prs.slide_masters[0].slide_layouts)

    footer_text = outline.get('footer', '')

    # Build title slide
    build_title_slide(prs, layouts, outline)

    # Build content slides
    for i, slide_def in enumerate(outline['slides']):
        slide_num = i + 2  # title is slide 1
        build_content_slide(prs, layouts, slide_def, slide_num, footer_text)

    prs.save(output_path)
    print(f"Built {len(prs.slides)} slides → {output_path}")


def main():
    parser = argparse.ArgumentParser(description='BET Deck Builder')
    parser.add_argument('--template', required=True, help='Path to Master_Ergebnis.pptx')
    parser.add_argument('--outline', required=True, help='Path to outline JSON')
    parser.add_argument('--output', required=True, help='Output .pptx path')
    args = parser.parse_args()

    with open(args.outline) as f:
        outline = json.load(f)

    build(args.template, outline, args.output)


if __name__ == '__main__':
    main()
